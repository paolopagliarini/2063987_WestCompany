#!/usr/bin/env python3
"""
Mars Habitat Automation Platform - University Presentation
Created by: Deda Daniel, Macrì Francesco, Pagliarini Paolo, Torresi Flavio

Professional PowerPoint presentation with Apple/Microsoft-style design
Following Elite PowerPoint Designer guidelines for world-class presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ============================================================================
# DESIGN SYSTEM - Tech Keynote Style (Apple/Tesla inspired)
# ============================================================================
COLORS = {
    'background': RGBColor(0x0A, 0x0A, 0x0A),      # Deep dark (#0A0A0A)
    'background_light': RGBColor(0x16, 0x16, 0x22), # Dark blue-gray
    'background_card': RGBColor(0x1E, 0x1E, 0x2E),  # Card background
    'primary': RGBColor(0xFF, 0xFF, 0xFF),          # White
    'secondary': RGBColor(0x94, 0x94, 0xB0),        # Light gray-blue
    'muted': RGBColor(0x6B, 0x6B, 0x80),            # Muted text
    'accent': RGBColor(0x00, 0xD4, 0xFF),           # Cyan blue (primary accent)
    'accent_alt': RGBColor(0x8B, 0x5C, 0xF6),       # Purple (secondary accent)
    'success': RGBColor(0x10, 0xB9, 0x81),          # Green
    'warning': RGBColor(0xF5, 0x9E, 0x0B),          # Orange
    'error': RGBColor(0xEF, 0x44, 0x44),            # Red
    'gradient_start': RGBColor(0x1A, 0x1A, 0x2E),
    'gradient_end': RGBColor(0x0A, 0x0A, 0x0A),
}

# Typography scale (following design system)
TYPOGRAPHY = {
    'hero_title': 72,       # Main presentation titles
    'section_title': 54,    # Section headers
    'slide_title': 44,      # Individual slide titles
    'subtitle': 32,         # Subtitles
    'body_large': 26,       # Emphasized body text
    'body': 22,             # Regular body text
    'body_small': 18,       # Smaller body text
    'caption': 16,          # Captions and notes
    'micro': 14,            # Very small text
}

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def add_background(slide, color=COLORS['background']):
    """Add solid background color to slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_text_box(slide, left, top, width, height, text, font_size=22,
                 font_color=COLORS['primary'], bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Segoe UI', italic=False):
    """Add styled text box"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                        Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return textbox


def add_accent_bar(slide, left, top, width=0.12, height=0.8, color=None):
    """Add accent color bar for visual hierarchy"""
    if color is None:
        color = COLORS['accent']
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_card(slide, left, top, width, height, border_color=None, fill_color=None):
    """Add a rounded rectangle card with optional border accent"""
    if fill_color is None:
        fill_color = COLORS['background_card']

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(2)
    else:
        card.line.fill.background()
    return card


def add_icon_circle(slide, left, top, size=0.4, color=None):
    """Add a colored circle (icon placeholder)"""
    if color is None:
        color = COLORS['accent']
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top),
        Inches(size), Inches(size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    return circle


def add_bullet_point(slide, left, top, text, font_size=22, color=None):
    """Add a bullet point with text"""
    if color is None:
        color = COLORS['accent']
    bullet = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top + 0.12),
        Inches(0.1), Inches(0.1)
    )
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = color
    bullet.line.fill.background()

    text_box = add_text_box(slide, left + 0.25, top, 10, 0.5, text,
                            font_size=font_size, font_color=COLORS['secondary'])
    return text_box


def add_slide_number(slide, number, total):
    """Add slide number in bottom right"""
    add_text_box(slide, 12.2, 6.9, 0.8, 0.3,
                 f"{number}/{total}", font_size=TYPOGRAPHY['micro'],
                 font_color=COLORS['muted'], alignment=PP_ALIGN.RIGHT)


def add_footer_line(slide):
    """Add subtle footer line"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(7.1),
        Inches(11.7), Inches(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['background_light']
    line.line.fill.background()


TOTAL_SLIDES = 24

# ============================================================================
# SLIDE 1: Title Slide (Hero)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

# Decorative gradient line at top
top_line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(13.333), Inches(0.05)
)
top_line.fill.solid()
top_line.fill.fore_color.rgb = COLORS['accent']
top_line.line.fill.background()

# Main title
add_text_box(slide, 0.8, 2.0, 11.7, 1.5,
             "Mars Habitat Automation Platform",
             font_size=TYPOGRAPHY['hero_title'], font_color=COLORS['primary'],
             bold=True, alignment=PP_ALIGN.CENTER, font_name='Segoe UI Light')

# Subtitle
add_text_box(slide, 0.8, 3.5, 11.7, 0.8,
             "Distributed Event-Driven System for Habitat Monitoring & Control",
             font_size=TYPOGRAPHY['subtitle'], font_color=COLORS['secondary'],
             alignment=PP_ALIGN.CENTER)

# Accent line
accent_line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(5.0), Inches(4.4),
    Inches(3.3), Inches(0.04)
)
accent_line.fill.solid()
accent_line.fill.fore_color.rgb = COLORS['accent']
accent_line.line.fill.background()

# Authors
add_text_box(slide, 0.8, 5.3, 11.7, 0.6,
             "Deda Daniel  •  Macrì Francesco  •  Pagliarini Paolo  •  Torresi Flavio",
             font_size=TYPOGRAPHY['body'], font_color=COLORS['secondary'],
             alignment=PP_ALIGN.CENTER)

# Year and course
add_text_box(slide, 0.8, 6.3, 11.7, 0.4,
             "Laboratory of Advanced Programming – Hackathon Exam 2024/2025",
             font_size=TYPOGRAPHY['caption'], font_color=COLORS['muted'],
             alignment=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 2: Agenda
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.5, 11.7, 0.9,
             "Agenda",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.35, 0.12, 0.7)

agenda_items = [
    ("01", "Project Overview", "Problem statement and objectives"),
    ("02", "System Architecture", "Microservices and event-driven design"),
    ("03", "Data Ingestion", "REST polling and SSE streaming"),
    ("04", "Automation Engine", "IF-THEN rule processing"),
    ("05", "Frontend Application", "Real-time dashboard features"),
    ("06", "Technology Stack", "Tools and frameworks used"),
    ("07", "Conclusions", "Achievements and future work"),
]

y_pos = 1.8
for num, title, desc in agenda_items:
    # Number
    add_text_box(slide, 0.8, y_pos, 0.7, 0.5,
                 num, font_size=TYPOGRAPHY['body_large'],
                 font_color=COLORS['accent'], bold=True)
    # Title
    add_text_box(slide, 1.6, y_pos, 4, 0.4,
                 title, font_size=TYPOGRAPHY['body'],
                 font_color=COLORS['primary'], bold=True)
    # Description
    add_text_box(slide, 1.6, y_pos + 0.35, 8, 0.3,
                 desc, font_size=TYPOGRAPHY['caption'],
                 font_color=COLORS['muted'])
    y_pos += 0.75

add_slide_number(slide, 2, TOTAL_SLIDES)

# ============================================================================
# SLIDE 3: Problem Statement
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.5, 11.7, 0.9,
             "The Challenge",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.35, 0.12, 0.7)

# Challenge description
add_text_box(slide, 0.8, 1.9, 11.7, 1.0,
             "A Mars habitat's automation stack has partially failed, leaving devices operating with incompatible protocols.",
             font_size=TYPOGRAPHY['body_large'], font_color=COLORS['secondary'])

# Challenge cards
challenges = [
    ("Multiple Data Sources", "REST sensors, SSE streams, WebSocket telemetry"),
    ("Heterogeneous Formats", "Different payload schemas requiring normalization"),
    ("Real-time Requirements", "Immediate response to environmental changes"),
    ("Automation Needs", "Rules-based actuator control for habitat safety"),
]

positions = [(0.8, 3.2), (6.8, 3.2), (0.8, 5.1), (6.8, 5.1)]

for i, ((title, desc), (left, top)) in enumerate(zip(challenges, positions)):
    card = add_card(slide, left, top, 5.7, 1.7, border_color=COLORS['accent'])

    # Icon
    add_icon_circle(slide, left + 0.3, top + 0.3, 0.35, COLORS['accent'])

    # Title
    add_text_box(slide, left + 0.8, top + 0.25, 4.5, 0.5,
                 title, font_size=TYPOGRAPHY['body'],
                 font_color=COLORS['primary'], bold=True)

    # Description
    add_text_box(slide, left + 0.3, top + 0.9, 5.1, 0.7,
                 desc, font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['secondary'])

add_slide_number(slide, 3, TOTAL_SLIDES)

# ============================================================================
# SLIDE 4: Solution Overview
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.5, 11.7, 0.9,
             "Our Solution",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.35, 0.12, 0.7)

add_text_box(slide, 0.8, 1.9, 11.7, 0.7,
             "A distributed event-driven platform with unified monitoring and control",
             font_size=TYPOGRAPHY['body_large'], font_color=COLORS['secondary'])

# Solution pillars
pillars = [
    ("Unified Ingestion", "REST polling + SSE streams normalized to single event schema", COLORS['accent']),
    ("Event-Driven Core", "RabbitMQ message broker for decoupled microservices", COLORS['accent_alt']),
    ("Smart Automation", "IF-THEN rules with automatic actuator control", COLORS['success']),
    ("Real-Time Dashboard", "React SPA with live updates and manual controls", COLORS['warning']),
]

y_pos = 2.7
for title, desc, color in pillars:
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(y_pos + 0.15),
        Inches(0.06), Inches(0.4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

    add_text_box(slide, 1.1, y_pos, 4, 0.4,
                 title, font_size=TYPOGRAPHY['body'],
                 font_color=COLORS['primary'], bold=True)

    add_text_box(slide, 1.1, y_pos + 0.4, 10, 0.4,
                 desc, font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['muted'])
    y_pos += 0.95

add_slide_number(slide, 4, TOTAL_SLIDES)

# ============================================================================
# SLIDE 5: System Architecture (with diagram)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "System Architecture",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# Add architecture diagram
diagram_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "diagrams", "asset/Diagram.png")
if os.path.exists(diagram_path):
    slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.5), width=Inches(12.3))

add_slide_number(slide, 5, TOTAL_SLIDES)

# ============================================================================
# SLIDE 6: Microservices Overview
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "Microservices Architecture",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# Services grid
services = [
    ("Ingestion", "8001", "REST polling & SSE streams\nSchema normalization", COLORS['accent']),
    ("Automation Engine", "8002", "Rule evaluation\nEvent-driven triggers", COLORS['accent_alt']),
    ("Rule Manager", "8003", "CRUD API for automation\nrules persistence", COLORS['success']),
    ("Notification", "8004", "Real-time SSE alerts\nEvent broadcasting", COLORS['warning']),
    ("Actuator Control", "8005", "Device command execution\nAudit logging", COLORS['error']),
    ("Data History", "8006", "Historical queries\nIn-memory cache", COLORS['accent']),
]

positions = [(0.8, 1.5), (4.6, 1.5), (8.4, 1.5), (0.8, 4.2), (4.6, 4.2), (8.4, 4.2)]

for (name, port, desc, color), (left, top) in zip(services, positions):
    # Service card
    card = add_card(slide, left, top, 3.5, 2.5, border_color=color)

    # Port badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + 2.5), Inches(top + 0.15),
        Inches(0.85), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_text_box(slide, left + 2.5, top + 0.12, 0.85, 0.35,
                 port, font_size=TYPOGRAPHY['micro'], font_color=COLORS['background'],
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Service name
    add_text_box(slide, left + 0.2, top + 0.2, 2.2, 0.5,
                 name, font_size=TYPOGRAPHY['body'],
                 font_color=COLORS['primary'], bold=True)

    # Description
    add_text_box(slide, left + 0.2, top + 0.8, 3.1, 1.5,
                 desc, font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['secondary'])

add_slide_number(slide, 6, TOTAL_SLIDES)

# ============================================================================
# SLIDE 7: Infrastructure Services
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "Infrastructure Services",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# RabbitMQ Box
rabbit_box = add_card(slide, 0.8, 1.5, 5.5, 2.3, border_color=COLORS['warning'])

add_text_box(slide, 1.0, 1.7, 5, 0.5,
             "RabbitMQ Message Broker", font_size=TYPOGRAPHY['body_large'],
             font_color=COLORS['warning'], bold=True)

add_text_box(slide, 1.0, 2.3, 5.1, 1.3,
             "• AMQP protocol (port 5672)\n• Topic exchange: mars_events\n• Event-driven pub/sub architecture\n• Durable queues with TTL\n• Management UI on port 15672",
             font_size=TYPOGRAPHY['body_small'], font_color=COLORS['secondary'])

# PostgreSQL Box
postgres_box = add_card(slide, 7, 1.5, 5.5, 2.3, border_color=COLORS['accent_alt'])

add_text_box(slide, 7.2, 1.7, 5, 0.5,
             "PostgreSQL Database", font_size=TYPOGRAPHY['body_large'],
             font_color=COLORS['accent_alt'], bold=True)

add_text_box(slide, 7.2, 2.3, 5.1, 1.3,
             "• PostgreSQL 16 (port 5433)\n• automation_rules table\n• sensor_readings history\n• actuator_commands audit\n• Persistent volume storage",
             font_size=TYPOGRAPHY['body_small'], font_color=COLORS['secondary'])

# Event Schema Box
schema_box = add_card(slide, 0.8, 4.1, 11.7, 2.5, border_color=COLORS['accent'])

add_text_box(slide, 1.0, 4.3, 11.3, 0.5,
             "Unified Event Schema", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['accent'], bold=True)

add_text_box(slide, 1.0, 4.85, 11.3, 0.6,
             '{ event_id, sensor_id, timestamp, metric, value, unit, source, status, raw_schema }',
             font_size=TYPOGRAPHY['body_small'], font_color=COLORS['primary'],
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.0, 5.55, 11.3, 0.8,
             "All sensor data normalized to a single format • REST polls + SSE streams → unified events",
             font_size=TYPOGRAPHY['caption'], font_color=COLORS['muted'],
             alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 7, TOTAL_SLIDES)

# ============================================================================
# SLIDE 8: Data Ingestion
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "Data Ingestion Service",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# REST Sensors Section
add_text_box(slide, 0.8, 1.5, 5.5, 0.5,
             "REST Sensors (Polling)", font_size=TYPOGRAPHY['body_large'],
             font_color=COLORS['accent'], bold=True)

rest_sensors = [
    ("greenhouse_temperature", "temperature"),
    ("entrance_humidity", "humidity"),
    ("co2_hall", "gas"),
    ("hydroponic_ph", "chemistry"),
    ("water_tank_level", "level"),
    ("corridor_pressure", "pressure"),
    ("air_quality_pm25", "particulate"),
    ("air_quality_voc", "chemistry"),
]

y_pos = 2.1
for sensor, stype in rest_sensors:
    add_text_box(slide, 1.0, y_pos, 4, 0.35,
                 f"• {sensor}", font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['secondary'])
    y_pos += 0.32

# Telemetry Streams Section
add_text_box(slide, 7, 1.5, 5.5, 0.5,
             "Telemetry Streams (SSE)", font_size=TYPOGRAPHY['body_large'],
             font_color=COLORS['accent_alt'], bold=True)

telemetry = [
    ("mars/telemetry/solar_array", "power"),
    ("mars/telemetry/radiation", "environment"),
    ("mars/telemetry/life_support", "environment"),
    ("mars/telemetry/thermal_loop", "thermal"),
    ("mars/telemetry/power_bus", "power"),
    ("mars/telemetry/power_consumption", "power"),
    ("mars/telemetry/airlock", "airlock"),
]

y_pos = 2.1
for topic, ttype in telemetry:
    add_text_box(slide, 7.2, y_pos, 5, 0.35,
                 f"• {topic.split('/')[-1]}", font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['secondary'])
    y_pos += 0.32

# Normalization note
note_box = add_card(slide, 0.8, 4.9, 11.7, 1.8, border_color=COLORS['success'])

add_text_box(slide, 1.0, 5.1, 11.3, 0.5,
             "Normalization Pipeline", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['success'], bold=True)

add_text_box(slide, 1.0, 5.6, 11.3, 0.9,
             "8 different REST schemas + 7 telemetry schemas → Single unified event format\nPublished to RabbitMQ with routing key: events.sensor.{sensor_id}",
             font_size=TYPOGRAPHY['body_small'], font_color=COLORS['secondary'])

add_slide_number(slide, 8, TOTAL_SLIDES)

# ============================================================================
# SLIDE 9: Raw Data to JSON (Normalization)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "Raw Data to JSON (Normalization)",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

add_text_box(slide, 0.8, 1.6, 11.7, 0.5,
             "Standardizing Telemetry and Events",
             font_size=TYPOGRAPHY['body_large'], font_color=COLORS['accent'], bold=True)

# Example box 1: Raw payload
add_card(slide, 0.8, 2.3, 5.0, 4.0, border_color=COLORS['error'])
add_text_box(slide, 1.0, 2.5, 4.6, 0.4, "Raw Telemetry (Varies)", font_size=TYPOGRAPHY['body'], font_color=COLORS['error'], bold=True)
raw_json = '''{
  "sensor": "DHT22",
  "data": {
    "temp": 39.99,
    "flow": 12.0
  }
}'''
add_text_box(slide, 1.0, 3.1, 4.6, 2.5, raw_json, font_size=TYPOGRAPHY['body_small'], font_color=COLORS['secondary'], font_name='Courier New')

# Arrow pointing right
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.1), Inches(3.8), Inches(0.8), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = COLORS['accent_alt']

# Example box 2: Unified format
add_card(slide, 7.2, 2.3, 5.3, 4.0, border_color=COLORS['success'])
add_text_box(slide, 7.4, 2.5, 4.6, 0.4, "Unified Event JSON", font_size=TYPOGRAPHY['body'], font_color=COLORS['success'], bold=True)
unified_json = '''{
  "event_id": "evt_101",
  "sensor_id": "thermal_loop",
  "timestamp": "2026-03-10T12:00:00Z",
  "metric": "temperature_c",
  "value": 39.99,
  "unit": "C",
  "source": "sse",
  "status": "warning"
}'''
add_text_box(slide, 7.4, 3.1, 4.8, 3.2, unified_json, font_size=TYPOGRAPHY['body_small'], font_color=COLORS['primary'], font_name='Courier New')

add_slide_number(slide, 9, TOTAL_SLIDES)

# ============================================================================
# SLIDE 10: In-Memory Cache & Event Mapping
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "In-Memory Cache & Mapping",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# The cache explanation
add_text_box(slide, 0.8, 1.5, 11.7, 0.8,
             "Ingestion service maintains the latest values for every sensor, exposed via GET /sensors/latest\nto the frontend with real-time updates.",
             font_size=TYPOGRAPHY['body_large'], font_color=COLORS['primary'])

# Mapping Details
add_card(slide, 0.8, 2.6, 11.7, 4.2, border_color=COLORS['accent'])
add_text_box(slide, 1.2, 2.8, 10.9, 0.5, "Event Mapping Strategies", font_size=TYPOGRAPHY['body'], font_color=COLORS['accent'], bold=True)

mappings = [
    "rest.scalar.v1 -> 1 direct event (1:1 mapping)",
    "rest.chemistry.v1 -> 1 event per measurement in the measurements array",
    "rest.particulate.v1 -> 3 separated events (PM1, PM2.5, PM10)",
    "rest.level.v1 -> 2 events (level_pct in % and level_liters in L)",
    "topic.power.v1 / environment.v1 -> 1 event per metric",
    "topic.thermal_loop.v1 -> 2 events (temperature_c, flow_l_min)",
    "topic.airlock.v1 -> 1 event (last_state encoded as status)"
]

y_pos = 3.5
for m in mappings:
    parts = m.split("->")
    if len(parts) == 2:
        schema = parts[0].strip()
        mapping = parts[1].strip()
        # draw small rectangle for schema
        add_text_box(slide, 1.2, y_pos, 4.3, 0.4, schema, font_size=TYPOGRAPHY['body_small'], font_color=COLORS['warning'], font_name='Courier New')
        # arrow symbol ->
        add_text_box(slide, 5.6, y_pos, 0.5, 0.4, "→", font_size=TYPOGRAPHY['body_small'], font_color=COLORS['secondary'])
        # mapping text
        add_text_box(slide, 6.1, y_pos, 6.0, 0.4, mapping, font_size=TYPOGRAPHY['body_small'], font_color=COLORS['secondary'])
    y_pos += 0.43

add_slide_number(slide, 10, TOTAL_SLIDES)

# ============================================================================
# SLIDE 11: Automation Engine
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "Automation Engine",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# Flow diagram
add_text_box(slide, 0.8, 1.4, 11.7, 0.5,
             "Event-Driven Rule Processing",
             font_size=TYPOGRAPHY['body_large'], font_color=COLORS['accent'], bold=True)

# Flow steps with arrows
steps = [
    ("1", "Receive", "RabbitMQ\nevent"),
    ("2", "Parse", "Extract\nsensor_id"),
    ("3", "Match", "Find rules"),
    ("4", "Evaluate", "Check\ncondition"),
    ("5", "Execute", "Publish\ncommand"),
]

x_positions = [0.8, 3.3, 5.8, 8.3, 10.8]
for i, (num, title, desc) in enumerate(steps):
    left = x_positions[i]

    # Step circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + 0.35), Inches(2.1),
        Inches(0.6), Inches(0.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['accent']
    circle.line.fill.background()
    add_text_box(slide, left + 0.35, Inches(2.15).inches, 0.6, 0.5,
                 num, font_size=TYPOGRAPHY['body'], font_color=COLORS['background'],
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, left, 2.8, 1.8, 0.4,
                 title, font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['primary'], bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, left, 3.2, 1.8, 0.7,
                 desc, font_size=TYPOGRAPHY['caption'],
                 font_color=COLORS['muted'], alignment=PP_ALIGN.CENTER)

    # Arrow (except last)
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(left + 1.85), Inches(2.3),
            Inches(0.4), Inches(0.25)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLORS['accent_alt']
        arrow.line.fill.background()

# Rule syntax box
rule_box = add_card(slide, 0.8, 4.3, 11.7, 2.4, border_color=COLORS['success'])

add_text_box(slide, 1.0, 4.5, 11.3, 0.5,
             "IF-THEN Rule Model", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['success'], bold=True)

add_text_box(slide, 1.0, 5.0, 11.3, 0.6,
             "IF <sensor_id> <operator> <value> THEN set <actuator_id> to ON | OFF",
             font_size=TYPOGRAPHY['body'], font_color=COLORS['primary'],
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1.0, 5.7, 11.3, 0.8,
             "Operators:  <   <=   =   >   >=   •   Example: IF greenhouse_temperature > 28 THEN set cooling_fan to ON",
             font_size=TYPOGRAPHY['body_small'], font_color=COLORS['secondary'],
             alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 11, TOTAL_SLIDES)

# ============================================================================
# SLIDE 12: Frontend Overview
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "Frontend Application",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# Tech stack
add_text_box(slide, 0.8, 1.4, 4, 0.5,
             "Technology Stack", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['accent'], bold=True)

tech_stack = [
    "React 18 + TypeScript",
    "Vite 6 Build Tool",
    "React Router 7",
    "Tailwind CSS 4",
    "Recharts for Charts",
    "Radix UI Components",
]

y_pos = 1.9
for tech in tech_stack:
    add_text_box(slide, 1.0, y_pos, 4, 0.35,
                 f"• {tech}", font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['secondary'])
    y_pos += 0.35

# Pages overview
add_text_box(slide, 5.3, 1.4, 7, 0.5,
             "Application Pages", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['accent'], bold=True)

pages = [
    ("Sensors", "REST sensors dashboard with real-time updates"),
    ("Telemetry", "Streaming data visualization with charts"),
    ("Actuators", "Device control and command history"),
    ("Rule Builder", "Create and edit automation rules"),
    ("Rules", "List, enable/disable, delete rules"),
    ("Status", "System health and notification log"),
]

y_pos = 1.9
for page, desc in pages:
    add_text_box(slide, 5.5, y_pos, 1.6, 0.35,
                 page, font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['primary'], bold=True)
    add_text_box(slide, 7.1, y_pos, 5, 0.35,
                 desc, font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['secondary'])
    y_pos += 0.45

# Real-time features
features_box = add_card(slide, 0.8, 4.7, 11.7, 2.0, border_color=COLORS['accent'])

add_text_box(slide, 1.0, 4.9, 11.3, 0.5,
             "Real-Time Features", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['accent'], bold=True)

add_text_box(slide, 1.0, 5.4, 11.3, 1.0,
             "Server-Sent Events (SSE) for live notifications  •  Aggressive polling for sensor updates\nReact state management for instant UI feedback  •  Recharts for live data visualization",
             font_size=TYPOGRAPHY['body_small'], font_color=COLORS['secondary'],
             alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 12, TOTAL_SLIDES)

# ============================================================================
# SLIDE 13: Sensors Dashboard (Mockup)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.3, 11.7, 0.6,
             "Sensors Dashboard",
             font_size=TYPOGRAPHY['slide_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 0.85, 0.1, 0.4)

add_text_box(slide, 2, 0.3, 10, 0.5,
             "Real-time REST sensor monitoring with status indicators",
             font_size=TYPOGRAPHY['caption'], font_color=COLORS['muted'])

# Add mockup image
mockup_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "mockups", "01 Sensors.png")
if os.path.exists(mockup_path):
    slide.shapes.add_picture(mockup_path, Inches(0.5), Inches(1.2), width=Inches(12.3))

add_slide_number(slide, 13, TOTAL_SLIDES)

# ============================================================================
# SLIDE 14: Telemetry Page (Mockup)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.3, 11.7, 0.6,
             "Telemetry Streaming",
             font_size=TYPOGRAPHY['slide_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 0.85, 0.1, 0.4)

add_text_box(slide, 2.8, 0.3, 10, 0.5,
             "SSE-powered live data visualization",
             font_size=TYPOGRAPHY['caption'], font_color=COLORS['muted'])

mockup_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "mockups", "02 Telemetry.png")
if os.path.exists(mockup_path):
    slide.shapes.add_picture(mockup_path, Inches(0.5), Inches(1.2), width=Inches(12.3))

add_slide_number(slide, 14, TOTAL_SLIDES)

# ============================================================================
# SLIDE 15: Actuators Page (Mockup)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.3, 11.7, 0.6,
             "Actuator Control",
             font_size=TYPOGRAPHY['slide_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 0.85, 0.1, 0.4)

add_text_box(slide, 2.5, 0.3, 10, 0.5,
             "Manual control and command history audit trail",
             font_size=TYPOGRAPHY['caption'], font_color=COLORS['muted'])

mockup_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "mockups", "03 Actuators.png")
if os.path.exists(mockup_path):
    slide.shapes.add_picture(mockup_path, Inches(0.5), Inches(1.2), width=Inches(12.3))

add_slide_number(slide, 15, TOTAL_SLIDES)

# ============================================================================
# SLIDE 16: Rule Builder Page (Mockup)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.3, 11.7, 0.6,
             "Rule Builder",
             font_size=TYPOGRAPHY['slide_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 0.85, 0.1, 0.4)

add_text_box(slide, 2.3, 0.3, 10, 0.5,
             "Create IF-THEN automation rules visually",
             font_size=TYPOGRAPHY['caption'], font_color=COLORS['muted'])

mockup_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "mockups", "04 rule Builder.png")
if os.path.exists(mockup_path):
    slide.shapes.add_picture(mockup_path, Inches(0.5), Inches(1.2), width=Inches(12.3))

add_slide_number(slide, 16, TOTAL_SLIDES)

# ============================================================================
# SLIDE 17: Rules Page (Mockup)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.3, 11.7, 0.6,
             "Rules Management",
             font_size=TYPOGRAPHY['slide_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 0.85, 0.1, 0.4)

add_text_box(slide, 2.8, 0.3, 10, 0.5,
             "List, toggle, edit, and delete automation rules",
             font_size=TYPOGRAPHY['caption'], font_color=COLORS['muted'])

mockup_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "mockups", "05 Rules.png")
if os.path.exists(mockup_path):
    slide.shapes.add_picture(mockup_path, Inches(0.5), Inches(1.2), width=Inches(12.3))

add_slide_number(slide, 17, TOTAL_SLIDES)

# ============================================================================
# SLIDE 18: Status Page (Mockup)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.3, 11.7, 0.6,
             "System Status",
             font_size=TYPOGRAPHY['slide_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 0.85, 0.1, 0.4)

add_text_box(slide, 2.3, 0.3, 10, 0.5,
             "Health monitoring and notification center",
             font_size=TYPOGRAPHY['caption'], font_color=COLORS['muted'])

mockup_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "mockups", "06 Status.png")
if os.path.exists(mockup_path):
    slide.shapes.add_picture(mockup_path, Inches(0.5), Inches(1.2), width=Inches(12.3))

add_slide_number(slide, 18, TOTAL_SLIDES)

# ============================================================================
# SLIDE 19: Database Schema
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "Database Design",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# Three tables
tables = [
    ("automation_rules", COLORS['accent'], [
        "id (SERIAL PK)",
        "name, description",
        "sensor_id, operator",
        "threshold_value, threshold_unit",
        "actuator_id, actuator_action",
        "is_active (BOOLEAN)",
        "created_at, updated_at"
    ]),
    ("sensor_readings", COLORS['accent_alt'], [
        "id (BIGSERIAL PK)",
        "sensor_id (VARCHAR)",
        "value (DECIMAL)",
        "unit (VARCHAR)",
        "source (rest/stream)",
        "recorded_at (TIMESTAMP)",
        "created_at (TIMESTAMP)"
    ]),
    ("actuator_commands", COLORS['warning'], [
        "id (BIGSERIAL PK)",
        "actuator_id (VARCHAR)",
        "previous_state, new_state",
        "source (auto/manual)",
        "reason (TEXT)",
        "rule_id (FK)",
        "executed_at (TIMESTAMP)"
    ]),
]

x_positions = [0.5, 4.6, 8.7]
for i, (name, color, fields) in enumerate(tables):
    left = x_positions[i]

    # Table box
    box = add_card(slide, left, 1.5, 4.0, 5.4, border_color=color)

    # Table name
    add_text_box(slide, left + 0.2, 1.7, 3.6, 0.5,
                 name, font_size=TYPOGRAPHY['body_small'], font_color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Separator line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + 0.3), Inches(2.3),
        Inches(3.4), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

    # Fields
    y_pos = 2.5
    for field in fields:
        add_text_box(slide, left + 0.3, y_pos, 3.4, 0.35,
                     f"• {field}", font_size=TYPOGRAPHY['caption'],
                     font_color=COLORS['secondary'])
        y_pos += 0.38

add_slide_number(slide, 19, TOTAL_SLIDES)

# ============================================================================
# SLIDE 20: Data Flow
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "Data Flow Architecture",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# Flow stages
stages = [
    ("IoT Simulator", "REST API\nSSE Streams", COLORS['warning']),
    ("Ingestion", "Normalize\nEvents", COLORS['accent']),
    ("RabbitMQ", "Route\nEvents", COLORS['warning']),
    ("Consumers", "Process &\nStore", COLORS['accent_alt']),
    ("Frontend", "Display &\nControl", COLORS['success']),
]

# Draw flow
y_center = 3.3
for i, (name, desc, color) in enumerate(stages):
    left = 0.6 + i * 2.5

    # Stage box
    box = add_card(slide, left, y_center - 0.8, 2.2, 1.6, border_color=color)

    # Name
    add_text_box(slide, left + 0.1, y_center - 0.6, 2, 0.5,
                 name, font_size=TYPOGRAPHY['body_small'], font_color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, left + 0.1, y_center + 0.1, 2, 0.7,
                 desc, font_size=TYPOGRAPHY['caption'], font_color=COLORS['secondary'],
                 alignment=PP_ALIGN.CENTER)

    # Arrow (except last)
    if i < len(stages) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(left + 2.25), Inches(y_center - 0.15),
            Inches(0.3), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLORS['accent']
        arrow.line.fill.background()

# Event types
add_text_box(slide, 0.8, 5.3, 11.7, 0.5,
             "Event Routing Keys", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['accent'], bold=True)

event_types = [
    ("Sensor Events", "events.sensor.{sensor_id}", COLORS['accent']),
    ("Actuator Commands", "commands.actuator.{actuator_id}", COLORS['accent_alt']),
    ("Notifications", "events.#", COLORS['warning']),
]

x_pos = 0.8
for name, pattern, color in event_types:
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_pos), Inches(5.85),
        Inches(0.06), Inches(0.4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

    add_text_box(slide, x_pos + 0.15, 5.8, 3.5, 0.4,
                 f"{name}:", font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['primary'], bold=True)
    add_text_box(slide, x_pos + 0.15, 6.2, 3.5, 0.4,
                 pattern, font_size=TYPOGRAPHY['caption'],
                 font_color=COLORS['muted'])
    x_pos += 4

add_slide_number(slide, 20, TOTAL_SLIDES)

# ============================================================================
# SLIDE 21: Technology Stack
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "Technology Stack",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# Backend
backend_box = add_card(slide, 0.5, 1.5, 3.8, 5.3, border_color=COLORS['accent'])

add_text_box(slide, 0.7, 1.7, 3.4, 0.5,
             "Backend Services", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['accent'], bold=True)

backend_tech = [
    "Python 3.11+",
    "FastAPI (Async)",
    "aio-pika (RabbitMQ)",
    "asyncpg + SQLAlchemy",
    "httpx (HTTP Client)",
    "Pydantic (Validation)",
    "Uvicorn (ASGI)",
]
y_pos = 2.4
for tech in backend_tech:
    add_text_box(slide, 0.9, y_pos, 3.2, 0.4,
                 f"• {tech}", font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['secondary'])
    y_pos += 0.4

# Frontend
frontend_box = add_card(slide, 4.75, 1.5, 3.8, 5.3, border_color=COLORS['accent_alt'])

add_text_box(slide, 4.95, 1.7, 3.4, 0.5,
             "Frontend", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['accent_alt'], bold=True)

frontend_tech = [
    "React 18",
    "TypeScript",
    "Vite 6",
    "React Router 7",
    "Tailwind CSS 4",
    "Recharts",
    "Radix UI",
]
y_pos = 2.4
for tech in frontend_tech:
    add_text_box(slide, 5.15, y_pos, 3.2, 0.4,
                 f"• {tech}", font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['secondary'])
    y_pos += 0.4

# Infrastructure
infra_box = add_card(slide, 9, 1.5, 3.8, 5.3, border_color=COLORS['warning'])

add_text_box(slide, 9.2, 1.7, 3.4, 0.5,
             "Infrastructure", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['warning'], bold=True)

infra_tech = [
    "PostgreSQL 16",
    "RabbitMQ (AMQP)",
    "Docker",
    "Docker Compose",
    "Nginx (Reverse Proxy)",
]
y_pos = 2.4
for tech in infra_tech:
    add_text_box(slide, 9.4, y_pos, 3.2, 0.4,
                 f"• {tech}", font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['secondary'])
    y_pos += 0.4

# Architecture note
add_text_box(slide, 9.2, 4.7, 3.4, 1.8,
             "Containerized\nMicroservices\nwith\nEvent-Driven\nArchitecture",
             font_size=TYPOGRAPHY['body_small'], font_color=COLORS['muted'],
             alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 21, TOTAL_SLIDES)

# ============================================================================
# SLIDE 22: User Stories Summary
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.4, 11.7, 0.7,
             "Key User Stories",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.05, 0.12, 0.5)

# Stories in categories
categories = [
    ("Monitoring", COLORS['accent'], [
        "Real-time sensor values",
        "Live telemetry from SSE",
        "Automatic dashboard updates",
        "Status indicators (ok/warning)",
        "Sensor filtering by category"
    ]),
    ("Control", COLORS['accent_alt'], [
        "View actuator ON/OFF states",
        "Manual actuator toggle",
        "Automatic rule-based control",
        "Command history audit"
    ]),
    ("Automation", COLORS['success'], [
        "Create IF-THEN rules",
        "View all automation rules",
        "Edit existing rules",
        "Delete/enable/disable",
        "Real-time trigger alerts"
    ]),
    ("Analytics", COLORS['warning'], [
        "Live time charts",
        "Historical data queries",
        "System health summary",
        "Connectivity monitoring"
    ]),
]

x_positions = [0.5, 3.7, 6.9, 10.1]
for i, (cat, color, stories) in enumerate(categories):
    left = x_positions[i]

    # Category box
    box = add_card(slide, left, 1.4, 2.9, 5.5, border_color=color)

    # Color accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(1.4),
        Inches(2.9), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    # Category name
    add_text_box(slide, left + 0.1, 1.6, 2.7, 0.5,
                 cat, font_size=TYPOGRAPHY['body'], font_color=color,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Stories
    y_pos = 2.3
    for story in stories:
        add_text_box(slide, left + 0.15, y_pos, 2.6, 0.7,
                     f"• {story}", font_size=TYPOGRAPHY['caption'],
                     font_color=COLORS['secondary'])
        y_pos += 0.65

add_slide_number(slide, 22, TOTAL_SLIDES)

# ============================================================================
# SLIDE 23: Conclusions
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

add_text_box(slide, 0.8, 0.5, 11.7, 0.8,
             "Conclusions",
             font_size=TYPOGRAPHY['section_title'], font_color=COLORS['primary'], bold=True)

add_accent_bar(slide, 0.8, 1.3, 0.12, 0.6)

# Achievements
add_text_box(slide, 0.8, 1.9, 11.7, 0.5,
             "Project Achievements", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['success'], bold=True)

achievements = [
    "Distributed microservices architecture with event-driven communication",
    "Real-time monitoring and control of simulated Mars habitat",
    "IF-THEN automation rules with automatic actuator control",
    "Modern React dashboard with SSE-based live updates",
    "Containerized deployment with Docker Compose",
    "Unified event schema for heterogeneous data sources",
]

y_pos = 2.4
for ach in achievements:
    bullet = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.8), Inches(y_pos + 0.12),
        Inches(0.1), Inches(0.1)
    )
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = COLORS['success']
    bullet.line.fill.background()

    add_text_box(slide, 1.1, y_pos, 11, 0.45,
                 ach, font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['secondary'])
    y_pos += 0.48

# Future Work
add_text_box(slide, 0.8, 5.4, 11.7, 0.5,
             "Future Enhancements", font_size=TYPOGRAPHY['body'],
             font_color=COLORS['accent_alt'], bold=True)

future = [
    ("Machine Learning", "Predictive maintenance and anomaly detection"),
    ("Mobile App", "Remote monitoring and push notifications"),
    ("Advanced Alerting", "Email/SMS notifications for critical events"),
]

y_pos = 5.9
for title, desc in future:
    add_text_box(slide, 1.0, y_pos, 2.5, 0.35,
                 f"→ {title}:", font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['primary'], bold=True)
    add_text_box(slide, 3.5, y_pos, 8, 0.35,
                 desc, font_size=TYPOGRAPHY['body_small'],
                 font_color=COLORS['muted'])
    y_pos += 0.4

add_slide_number(slide, 23, TOTAL_SLIDES)

# ============================================================================
# SLIDE 24: Thank You
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['background'])

# Decorative gradient line at top
top_line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(13.333), Inches(0.05)
)
top_line.fill.solid()
top_line.fill.fore_color.rgb = COLORS['accent']
top_line.line.fill.background()

# Thank you
add_text_box(slide, 0.8, 2.5, 11.7, 1.2,
             "Thank You",
             font_size=72, font_color=COLORS['primary'],
             bold=True, alignment=PP_ALIGN.CENTER, font_name='Segoe UI Light')

# Subtitle
add_text_box(slide, 0.8, 3.8, 11.7, 0.6,
             "Questions?",
             font_size=TYPOGRAPHY['subtitle'], font_color=COLORS['accent'],
             alignment=PP_ALIGN.CENTER)

# Accent line
accent_line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(5.5), Inches(4.5),
    Inches(2.3), Inches(0.04)
)
accent_line.fill.solid()
accent_line.fill.fore_color.rgb = COLORS['accent']
accent_line.line.fill.background()

# Team
add_text_box(slide, 0.8, 5.2, 11.7, 0.5,
             "Deda Daniel  •  Macrì Francesco  •  Pagliarini Paolo  •  Torresi Flavio",
             font_size=TYPOGRAPHY['body'], font_color=COLORS['secondary'],
             alignment=PP_ALIGN.CENTER)

# Course info
add_text_box(slide, 0.8, 6.3, 11.7, 0.4,
             "Laboratory of Advanced Programming – Hackathon Exam 2024/2025",
             font_size=TYPOGRAPHY['caption'], font_color=COLORS['muted'],
             alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 24, TOTAL_SLIDES)

# Save presentation
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Mars_Habitat_Presentation.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"🎨 Design: Tech Keynote Style (Apple/Tesla inspired)")
print(f"📐 Aspect ratio: 16:9")